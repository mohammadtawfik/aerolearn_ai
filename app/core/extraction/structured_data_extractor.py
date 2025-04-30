import os
from typing import List, Dict, Any, Union

try:
    import pdfplumber
except ImportError:
    pdfplumber = None
try:
    import docx
except ImportError:
    docx = None
try:
    import pptx
except ImportError:
    pptx = None

class StructuredDataExtractionError(Exception):
    pass

class StructuredDataExtractor:
    """
    Extracts tabular data from PDF, DOCX, and PPTX files.
    Only basic diagram detection is supported (future work can extend).
    """
    def extract_tables(self, filepath: str) -> List[List[Any]]:
        ext = os.path.splitext(filepath)[1].lower()
        if ext == '.pdf':
            return self._extract_tables_from_pdf(filepath)
        elif ext == '.docx':
            return self._extract_tables_from_docx(filepath)
        elif ext == '.pptx':
            return self._extract_tables_from_pptx(filepath)
        else:
            raise StructuredDataExtractionError(f"Unsupported file type: {ext}")

    def _extract_tables_from_pdf(self, filepath: str) -> List[List[Any]]:
        if not pdfplumber:
            raise ImportError("pdfplumber is required for PDF table extraction.")
        tables = []
        with pdfplumber.open(filepath) as pdf:
            for page in pdf.pages:
                page_tables = page.extract_tables()
                tables.extend(page_tables)
        return tables

    def _extract_tables_from_docx(self, filepath: str) -> List[List[Any]]:
        if not docx:
            raise ImportError("python-docx is required for DOCX table extraction.")
        doc = docx.Document(filepath)
        result = []
        for table in doc.tables:
            for row in table.rows:
                result.append([cell.text.strip() for cell in row.cells])
        return result

    def _extract_tables_from_pptx(self, filepath: str) -> List[List[Any]]:
        if not pptx:
            raise ImportError("python-pptx is required for PPTX table extraction.")
        prs = pptx.Presentation(filepath)
        tables = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_table:
                    tbl = shape.table
                    table_data = [
                        [cell.text for cell in row.cells]
                        for row in tbl.rows
                    ]
                    tables.extend(table_data)
        return tables

    def extract_diagrams(self, filepath: str) -> List[Any]:
        # Placeholder for diagram extraction; can employ OCR, layout-detection, or ML in future
        return []