import os
from typing import Union, List, Dict, Any

try:
    import pdfplumber
except ImportError:
    pdfplumber = None
try:
    import pytesseract
except ImportError:
    pytesseract = None
try:
    from PIL import Image
except ImportError:
    Image = None
try:
    import docx
except ImportError:
    docx = None
try:
    import pptx
except ImportError:
    pptx = None

class TextExtractionError(Exception):
    pass

class TextExtractor:
    """
    Extracts raw text from various document formats including PDF, DOCX, PPTX, and TXT.
    Falls back to OCR with tesseract or PIL if necessary for image-based content.
    
    Provides robust support for text files with various encodings and error handling.
    """
    def extract(self, filepath: str) -> str:
        """
        Main public API: extract text from a file, regardless of format.
        Delegates to extract_text().
        """
        return self.extract_text(filepath)
        
    def extract_text(self, filepath: str) -> str:
        ext = os.path.splitext(filepath)[1].lower()
        if ext == '.pdf':
            return self._extract_text_from_pdf(filepath)
        elif ext == '.docx':
            return self._extract_text_from_docx(filepath)
        elif ext == '.pptx':
            return self._extract_text_from_pptx(filepath)
        elif ext == '.txt':
            return self._extract_text_from_txt(filepath)
        else:
            raise TextExtractionError(f"Unsupported file type: {ext}")

    def _extract_text_from_pdf(self, filepath: str) -> str:
        if not pdfplumber:
            raise ImportError("pdfplumber is required for PDF text extraction.")
        text = ""
        with pdfplumber.open(filepath) as pdf:
            for page in pdf.pages:
                text += page.extract_text() or ""
                # If no text, try OCR
                if not text.strip() and pytesseract and Image:
                    im = page.to_image(resolution=300).original
                    text += pytesseract.image_to_string(im)
        return text

    def _extract_text_from_docx(self, filepath: str) -> str:
        if not docx:
            raise ImportError("python-docx is required for DOCX text extraction.")
        doc = docx.Document(filepath)
        return "\n".join(para.text for para in doc.paragraphs if para.text.strip())

    def _extract_text_from_pptx(self, filepath: str) -> str:
        if not pptx:
            raise ImportError("python-pptx is required for PPTX text extraction.")
        prs = pptx.Presentation(filepath)
        text_runs = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_runs.append(shape.text)
        return "\n".join(text_runs)
        
    def _extract_text_from_txt(self, filepath: str) -> str:
        """
        Extract plain text from a .txt file with robust encoding detection.
        
        Attempts to read the file with UTF-8 encoding first, then falls back
        to other common encodings if that fails. Handles potential errors
        during file reading and provides detailed error information.
        
        Args:
            filepath: Path to the text file
            
        Returns:
            Extracted text content as string
            
        Raises:
            TextExtractionError: If file cannot be read with any encoding
        """
        encodings = ['utf-8', 'latin-1', 'ascii', 'utf-16', 'windows-1252']
        
        for encoding in encodings:
            try:
                with open(filepath, "r", encoding=encoding) as f:
                    return f.read()
            except UnicodeDecodeError:
                # Try next encoding
                continue
            except FileNotFoundError:
                raise TextExtractionError(f"Text file not found: {filepath}")
            except PermissionError:
                raise TextExtractionError(f"Permission denied when accessing: {filepath}")
            except IsADirectoryError:
                raise TextExtractionError(f"Expected a file but got a directory: {filepath}")
            except Exception as e:
                # Try next encoding or raise if last one
                if encoding == encodings[-1]:
                    raise TextExtractionError(f"Failed to read TXT file with any encoding: {e}")
                
        # This should not be reached, but just in case
        raise TextExtractionError("Failed to read TXT file with all attempted encodings")
